"""Generate sample test documents (PDF, DOCX, PPTX) for the document viewer."""

import os

OUTPUT_DIR = os.path.join(os.path.dirname(__file__), "test_documents")


def generate_pdf():
    from reportlab.lib.pagesizes import letter
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import inch
    from reportlab.platypus import (
        SimpleDocTemplate,
        Paragraph,
        Spacer,
        Table,
        TableStyle,
        PageBreak,
    )
    from reportlab.lib import colors

    path = os.path.join(OUTPUT_DIR, "sample.pdf")
    doc = SimpleDocTemplate(path, pagesize=letter)
    styles = getSampleStyleSheet()
    story = []

    # Page 1
    story.append(Paragraph("Quarterly Business Report", styles["Title"]))
    story.append(Spacer(1, 0.3 * inch))
    story.append(Paragraph("Executive Summary", styles["Heading2"]))
    story.append(
        Paragraph(
            "This report provides an overview of the company's performance during Q3 2025. "
            "Revenue grew by 15% compared to the previous quarter, driven primarily by expansion "
            "into new markets and the successful launch of our cloud platform.",
            styles["BodyText"],
        )
    )
    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph("Key Metrics", styles["Heading2"]))

    table_data = [
        ["Metric", "Q2 2025", "Q3 2025", "Change"],
        ["Revenue ($M)", "42.3", "48.6", "+15%"],
        ["Active Users", "1.2M", "1.8M", "+50%"],
        ["Customer Satisfaction", "87%", "91%", "+4pp"],
        ["Employee Count", "320", "385", "+20%"],
    ]
    table = Table(table_data, colWidths=[2 * inch, 1.2 * inch, 1.2 * inch, 1 * inch])
    table.setStyle(
        TableStyle(
            [
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#0f62fe")),
                ("TEXTCOLOR", (0, 0), (-1, 0), colors.white),
                ("FONTNAME", (0, 0), (-1, 0), "Helvetica-Bold"),
                ("ALIGN", (1, 0), (-1, -1), "CENTER"),
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f4f4f4")]),
            ]
        )
    )
    story.append(table)

    # Page 2
    story.append(PageBreak())
    story.append(Paragraph("Market Analysis", styles["Heading1"]))
    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph("Regional Performance", styles["Heading2"]))
    story.append(
        Paragraph(
            "The North American market continues to be our strongest region, accounting for 60% "
            "of total revenue. European operations showed significant improvement with a 25% "
            "increase in new customer acquisition. The Asia-Pacific region remains our fastest "
            "growing market with 40% year-over-year growth.",
            styles["BodyText"],
        )
    )
    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph("Product Roadmap", styles["Heading2"]))
    story.append(
        Paragraph(
            "Looking ahead to Q4, we plan to release version 3.0 of our platform which includes "
            "AI-powered document analysis, enhanced collaboration features, and improved "
            "integration with third-party services. Beta testing is scheduled to begin in October.",
            styles["BodyText"],
        )
    )
    story.append(Spacer(1, 0.2 * inch))
    story.append(Paragraph("Conclusion", styles["Heading2"]))
    story.append(
        Paragraph(
            "The company is well-positioned for continued growth. Our strategic investments in "
            "technology and talent are yielding strong returns, and we remain confident in "
            "achieving our annual targets.",
            styles["BodyText"],
        )
    )

    doc.build(story)
    print(f"Generated: {path}")


def generate_docx():
    from docx import Document
    from docx.shared import Pt, Inches
    from docx.enum.text import WD_ALIGN_PARAGRAPH

    path = os.path.join(OUTPUT_DIR, "sample.docx")
    doc = Document()

    # Title
    title = doc.add_heading("Technical Design Document", level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.add_paragraph("Version 2.1 | January 2026", style="Subtitle")

    doc.add_heading("1. Introduction", level=1)
    doc.add_paragraph(
        "This document outlines the technical architecture for the Document Processing "
        "System (DPS). The system is designed to handle large-scale document ingestion, "
        "parsing, and analysis using modern cloud-native technologies."
    )

    doc.add_heading("2. System Architecture", level=1)
    doc.add_paragraph(
        "The DPS follows a microservices architecture deployed on Kubernetes. "
        "Each service is independently scalable and communicates via gRPC and "
        "asynchronous message queues."
    )

    doc.add_heading("2.1 Core Services", level=2)
    items = [
        "Document Ingestion Service - Handles file uploads and format validation",
        "Parser Service - Extracts text, tables, and metadata using Docling",
        "Analysis Service - Runs NLP models for entity extraction and classification",
        "Storage Service - Manages document storage in S3 and metadata in PostgreSQL",
        "API Gateway - Provides RESTful API access with authentication and rate limiting",
    ]
    for item in items:
        doc.add_paragraph(item, style="List Bullet")

    doc.add_heading("3. Data Flow", level=1)
    doc.add_paragraph(
        "Documents enter the system through the Ingestion Service, which validates "
        "the file format and stores the raw document in object storage. The Parser "
        "Service then processes the document asynchronously, extracting structured "
        "content. Results are stored in the database and made available through the API."
    )

    doc.add_heading("4. Security Considerations", level=1)
    doc.add_paragraph(
        "All data is encrypted at rest using AES-256 and in transit using TLS 1.3. "
        "Access control is managed through OAuth 2.0 with role-based permissions. "
        "Document content is isolated per tenant to ensure data privacy."
    )

    doc.add_heading("5. Performance Requirements", level=1)
    table = doc.add_table(rows=5, cols=3)
    table.style = "Light Grid Accent 1"
    headers = ["Metric", "Target", "Current"]
    for i, h in enumerate(headers):
        table.rows[0].cells[i].text = h
    data = [
        ("Document processing time", "< 5 seconds", "3.2 seconds"),
        ("API response time (p95)", "< 200ms", "145ms"),
        ("Throughput", "100 docs/min", "85 docs/min"),
        ("System uptime", "99.9%", "99.95%"),
    ]
    for row_idx, (m, t, c) in enumerate(data, 1):
        table.rows[row_idx].cells[0].text = m
        table.rows[row_idx].cells[1].text = t
        table.rows[row_idx].cells[2].text = c

    doc.save(path)
    print(f"Generated: {path}")


def generate_pptx():
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    path = os.path.join(OUTPUT_DIR, "sample.pptx")
    prs = Presentation()

    # Slide 1: Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Project Status Update"
    slide.placeholders[1].text = "Document Processing Initiative\nJanuary 2026"

    # Slide 2: Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Overview"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Goals for Q1 2026:"
    bullets = [
        "Launch document viewer with multi-format support",
        "Integrate Docling for intelligent document parsing",
        "Achieve 95% text extraction accuracy",
        "Support PDF, DOCX, and PPTX formats",
    ]
    for bullet in bullets:
        p = tf.add_paragraph()
        p.text = bullet
        p.level = 1

    # Slide 3: Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Technical Architecture"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "System Components:"
    components = [
        "React 18 Frontend with Carbon Design System",
        "Python FastAPI Backend",
        "Docling Document Processing Engine",
        "RESTful API with JSON responses",
    ]
    for comp in components:
        p = tf.add_paragraph()
        p.text = comp
        p.level = 1

    # Slide 4: Timeline
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Timeline & Milestones"
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = "Key milestones:"
    milestones = [
        "Week 1-2: Backend API and document processing",
        "Week 3-4: Frontend UI with three-panel layout",
        "Week 5: Integration testing and optimization",
        "Week 6: User acceptance testing and deployment",
    ]
    for m in milestones:
        p = tf.add_paragraph()
        p.text = m
        p.level = 1

    prs.save(path)
    print(f"Generated: {path}")


if __name__ == "__main__":
    os.makedirs(OUTPUT_DIR, exist_ok=True)
    generate_pdf()
    generate_docx()
    generate_pptx()
    print("All test documents generated successfully.")
